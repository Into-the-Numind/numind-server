"""
test_basic.py — Core unit tests for pptx-author skill.

Run from the skills/pptx-author directory:
    pytest tests/test_basic.py -v

Tests cover:
  - Slide count, 16:9 dimensions
  - Cover / bullet / body / table / chart / image / notes renderers
  - Bullet list truncation at 6
  - Chart types (bar, line, pie, scatter) via render_chart_to_image
  - Output file creation
  - Result warnings list structure
"""

from __future__ import annotations

import os
import sys
import tempfile

import pytest

# Make helpers importable from this test file
_skill_dir = os.path.join(os.path.dirname(__file__), "..")
sys.path.insert(0, _skill_dir)

from pptx import Presentation
from pptx.util import Cm
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture

from helpers import DeckBuilder, BrandConfig, merge_brand_config
from helpers.chart import render_chart_to_image
from helpers.brand import DEFAULT_BRAND_CONFIG


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture()
def default_brand():
    return BrandConfig.from_dict(None)


@pytest.fixture()
def custom_brand():
    return BrandConfig.from_dict({
        "primary_color": "#2563EB",
        "secondary_color": "#1E40AF",
        "company_name": "TestCo",
        "font_family": "Calibri",
    })


@pytest.fixture()
def tmp_output(tmp_path):
    return str(tmp_path / "output.pptx")


@pytest.fixture()
def sample_chart_def():
    return {
        "chart_type": "bar",
        "title": "Test Chart",
        "data": {
            "categories": ["A", "B", "C"],
            "series": [
                {"name": "Series 1", "values": [10, 20, 30]},
                {"name": "Series 2", "values": [15, 25, 35]},
            ],
        },
    }


# ---------------------------------------------------------------------------
# 1. Slide count
# ---------------------------------------------------------------------------

def test_basic_slide_creation(tmp_output, default_brand):
    """3 slides in → 3 slides in presentation."""
    builder = DeckBuilder(brand=default_brand)
    for layout in ("cover", "title-bullets", "end"):
        builder.add_slide({"layout": layout, "title": f"Test {layout}"})
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    assert len(prs.slides) == 3


# ---------------------------------------------------------------------------
# 2. 16:9 dimensions
# ---------------------------------------------------------------------------

def test_slide_16_9_dimensions(tmp_output, default_brand):
    """Slide width must be ≈ 33.87 cm (±0.1 cm tolerance)."""
    builder = DeckBuilder(brand=default_brand)
    builder.add_slide({"layout": "blank", "title": ""})
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    expected_w = Cm(33.87)
    assert abs(prs.slide_width - expected_w) < Cm(0.1), (
        f"Expected ~{expected_w}, got {prs.slide_width}"
    )


# ---------------------------------------------------------------------------
# 3. Cover title text
# ---------------------------------------------------------------------------

def test_cover_title_text(tmp_output, default_brand):
    """Cover slide must contain the title string in at least one textbox."""
    title = "My Cover Slide Title"
    builder = DeckBuilder(brand=default_brand)
    builder.add_slide({"layout": "cover", "title": title})
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    slide = prs.slides[0]
    all_text = " ".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )
    assert title in all_text


# ---------------------------------------------------------------------------
# 4. Bullet truncation to 6
# ---------------------------------------------------------------------------

def test_bullets_max_6(tmp_output, default_brand):
    """7 bullet_points provided → only 6 bullet lines rendered in textbox."""
    seven_bullets = [f"Bullet {i}" for i in range(1, 8)]
    builder = DeckBuilder(brand=default_brand)
    builder.add_slide({"layout": "title-bullets", "title": "Bullets", "bullet_points": seven_bullets})
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    slide = prs.slides[0]

    # Find the textbox that contains bullets (has '•' prefix)
    bullet_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text
                if t.startswith("•"):
                    bullet_texts.append(t)

    assert len(bullet_texts) <= 6, (
        f"Expected at most 6 bullet paragraphs, got {len(bullet_texts)}"
    )
    # The 7th bullet must not be present
    assert "Bullet 7" not in " ".join(bullet_texts)


# ---------------------------------------------------------------------------
# 5. Table shape present
# ---------------------------------------------------------------------------

def test_table_rendered(tmp_output, default_brand):
    """title-table layout must produce at least one Table shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    builder = DeckBuilder(brand=default_brand)
    builder.add_slide({
        "layout": "title-table",
        "title": "Test Table",
        "table": {
            "headers": ["Col A", "Col B"],
            "rows": [["r1c1", "r1c2"], ["r2c1", "r2c2"]],
        },
    })
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    slide = prs.slides[0]
    table_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
    assert len(table_shapes) >= 1


# ---------------------------------------------------------------------------
# 6. Chart image inserted (Picture shape)
# ---------------------------------------------------------------------------

def test_chart_image_inserted(tmp_output, default_brand):
    """title-chart layout must insert at least one Picture shape."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    builder = DeckBuilder(brand=default_brand)
    builder.add_slide({
        "layout": "title-chart",
        "title": "Chart Slide",
        "chart": {
            "chart_type": "bar",
            "title": "Revenue",
            "data": {
                "categories": ["Q1", "Q2"],
                "series": [{"name": "Rev", "values": [100, 200]}],
            },
        },
    })
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    slide = prs.slides[0]
    pic_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pic_shapes) >= 1


# ---------------------------------------------------------------------------
# 7. render_chart_to_image — bar, line, pie, scatter
# ---------------------------------------------------------------------------

def _simple_chart_def(chart_type: str, categories=None, series_values=None) -> dict:
    return {
        "chart_type": chart_type,
        "title": f"{chart_type} test",
        "data": {
            "categories": categories or ["X", "Y", "Z"],
            "series": [{"name": "S1", "values": series_values or [1, 2, 3]}],
        },
    }


def test_chart_bar_no_crash():
    buf = render_chart_to_image(_simple_chart_def("bar"))
    assert buf.read(4) == b"\x89PNG"  # valid PNG magic bytes


def test_chart_line_with_brand_color():
    """A red primary colour should produce red pixels in the PNG."""
    from PIL import Image
    import io

    buf = render_chart_to_image(
        _simple_chart_def("line"),
        primary_color="#FF0000",
    )
    img = Image.open(buf).convert("RGB")
    pixels = list(img.getdata())
    # Count pixels where R channel > 200 and G/B < 100 (strongly red)
    red_pixels = sum(1 for r, g, b in pixels if r > 200 and g < 100 and b < 100)
    assert red_pixels > 0, "Expected red pixels from brand primary colour"


def test_chart_pie_no_crash():
    buf = render_chart_to_image(
        _simple_chart_def("pie", ["Apple", "Banana", "Cherry"], [40, 35, 25])
    )
    assert len(buf.read()) > 1000  # non-trivial image size


def test_chart_scatter_no_crash():
    buf = render_chart_to_image(
        _simple_chart_def("scatter", ["1.0", "2.0", "3.0"], [10, 20, 15])
    )
    assert len(buf.read()) > 1000


def test_chart_scatter_string_categories_fallback():
    """scatter with non-numeric categories must not raise."""
    buf = render_chart_to_image(
        _simple_chart_def("scatter", ["Jan", "Feb", "Mar"], [10, 20, 15])
    )
    assert len(buf.read()) > 1000


def test_chart_invalid_type_raises():
    with pytest.raises(ValueError, match="Unsupported chart_type"):
        render_chart_to_image(_simple_chart_def("donut"))


# ---------------------------------------------------------------------------
# 8. Logo — missing file produces warning
# ---------------------------------------------------------------------------

def test_logo_missing_file_warning(tmp_output):
    """brand_config with non-existent logo_path must produce a warning."""
    brand = BrandConfig.from_dict({"logo_path": "nonexistent_logo.png"})
    builder = DeckBuilder(brand=brand)
    builder.add_slide({"layout": "cover", "title": "Test"})
    builder.apply_brand()
    builder.save(tmp_output)

    assert len(builder.warnings) > 0
    assert any("logo" in w.lower() for w in builder.warnings)


def test_svg_logo_warning(tmp_output):
    """SVG logo_path must trigger an explicit 'SVG not supported' warning."""
    brand = BrandConfig.from_dict({"logo_path": "logo.svg"})
    builder = DeckBuilder(brand=brand)
    builder.add_slide({"layout": "cover", "title": "Test"})
    builder.apply_brand()
    builder.save(tmp_output)

    assert any("svg" in w.lower() or "SVG" in w for w in builder.warnings)


# ---------------------------------------------------------------------------
# 9. Brand primary colour applied to cover title
# ---------------------------------------------------------------------------

def test_brand_primary_color_applied(tmp_output):
    """
    Cover slide with red primary_color must apply red in the presentation.

    The cover design uses primary_color as:
      (a) the background colour band (rectangle fill) — checked here
      (b) white text on top of the band — title text is WHITE by design

    Verify (a): the cover slide has at least one Shape whose fill
    fore-colour matches the primary_color (#FF0000).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor

    brand = BrandConfig.from_dict({"primary_color": "#FF0000"})  # red
    builder = DeckBuilder(brand=brand)
    builder.add_slide({"layout": "cover", "title": "Red Title"})
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    slide = prs.slides[0]

    found_primary_fill = False
    for shape in slide.shapes:
        try:
            fill = shape.fill
            if fill.type is not None:
                fore = fill.fore_color
                rgb = fore.rgb
                if rgb[0] > 200 and rgb[1] < 50 and rgb[2] < 50:
                    found_primary_fill = True
                    break
        except Exception:
            pass

    assert found_primary_fill, (
        "Expected at least one shape with red fill from primary_color=#FF0000. "
        "The cover uses primary_color as the header band background."
    )


# ---------------------------------------------------------------------------
# 10. Notes written to notes_slide
# ---------------------------------------------------------------------------

def test_notes_written(tmp_output, default_brand):
    """slide_def['notes'] must appear in notes_slide text."""
    note = "这是演讲者备注"
    builder = DeckBuilder(brand=default_brand)
    builder.add_slide({
        "layout": "title-body",
        "title": "Slide with Notes",
        "body": "Body text",
        "notes": note,
    })
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    slide = prs.slides[0]
    notes_text = slide.notes_slide.notes_text_frame.text
    assert note in notes_text


# ---------------------------------------------------------------------------
# 11. Output file exists with non-zero size
# ---------------------------------------------------------------------------

def test_output_file_created(tmp_output, default_brand):
    """Output .pptx file must exist and be non-empty after save."""
    builder = DeckBuilder(brand=default_brand)
    builder.add_slide({"layout": "title-bullets", "title": "Test", "bullet_points": ["a", "b"]})
    builder.save(tmp_output)

    assert os.path.exists(tmp_output)
    assert os.path.getsize(tmp_output) > 0


# ---------------------------------------------------------------------------
# 12. No brand logo → empty warnings list
# ---------------------------------------------------------------------------

def test_brand_applied_no_logo_no_warnings(tmp_output):
    """brand_config with no logo should produce zero warnings after apply_brand."""
    brand = BrandConfig.from_dict({"primary_color": "#2563EB", "logo_path": None})
    builder = DeckBuilder(brand=brand)
    builder.add_slide({"layout": "end", "title": "The End"})
    builder.apply_brand()
    builder.save(tmp_output)

    assert builder.warnings == []


# ---------------------------------------------------------------------------
# 13. merge_brand_config deep-merge behaviour
# ---------------------------------------------------------------------------

def test_merge_brand_config_partial_override():
    """Partial user config must not wipe unspecified default keys."""
    merged = merge_brand_config({"primary_color": "#123456"})
    assert merged["primary_color"] == "#123456"
    assert merged["secondary_color"] == DEFAULT_BRAND_CONFIG["secondary_color"]
    assert merged["company_name"] == DEFAULT_BRAND_CONFIG["company_name"]


def test_merge_brand_config_none_input():
    """None user_config must return exact DEFAULT_BRAND_CONFIG copy."""
    merged = merge_brand_config(None)
    for k, v in DEFAULT_BRAND_CONFIG.items():
        assert merged[k] == v


# ---------------------------------------------------------------------------
# 14. DeckBuilder: multi-layout smoke test (all 9 non-blank layouts)
# ---------------------------------------------------------------------------

def test_all_layouts_no_crash(tmp_output, default_brand):
    """Building one slide of each non-blank layout must not raise."""
    builder = DeckBuilder(brand=default_brand)
    layouts = [
        {"layout": "cover", "title": "Cover", "subtitle": "Sub"},
        {"layout": "section", "title": "Section", "subtitle": "01"},
        {"layout": "title-body", "title": "Body", "body": "Text"},
        {
            "layout": "title-bullets",
            "title": "Bullets",
            "bullet_points": ["A", "B"],
        },
        {
            "layout": "title-table",
            "title": "Table",
            "table": {"headers": ["H1", "H2"], "rows": [["a", "b"]]},
        },
        {
            "layout": "title-chart",
            "title": "Chart",
            "chart": {
                "chart_type": "line",
                "title": "L",
                "data": {
                    "categories": ["X", "Y"],
                    "series": [{"name": "S", "values": [1, 2]}],
                },
            },
        },
        {
            "layout": "title-image",
            "title": "Image",
            "image": {"path": "missing.jpg"},  # graceful missing-file handling
        },
        {
            "layout": "two-column",
            "title": "Two Col",
            "bullet_points": ["Left bullet"],
        },
        {"layout": "end", "title": "Done"},
    ]
    for sl in layouts:
        builder.add_slide(sl)
    builder.save(tmp_output)

    prs = Presentation(tmp_output)
    assert len(prs.slides) == len(layouts)
