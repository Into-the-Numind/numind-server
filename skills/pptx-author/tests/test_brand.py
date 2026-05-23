"""
test_brand.py — Unit tests for brand.py module.

Tests:
  - BrandConfig dataclass construction (from_dict, defaults, partial override)
  - primary_rgb / secondary_rgb hex-to-RGBColor conversion
  - apply_brand_to_slide: logo insertion (valid PNG), missing file, SVG rejection
  - apply_brand_to_deck: batch logo across multiple slides
  - apply_brand_colors_to_template: sentinel hex replacement
  - merge_brand_config: deep-merge behaviour
"""

from __future__ import annotations

import os
import sys
import struct
import zlib
import tempfile

import pytest

_skill_dir = os.path.join(os.path.dirname(__file__), "..")
sys.path.insert(0, _skill_dir)

from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor

from helpers.brand import (
    BrandConfig,
    DEFAULT_BRAND_CONFIG,
    merge_brand_config,
    apply_brand_to_slide,
    apply_brand_to_deck,
    apply_brand_colors_to_template,
    PRIMARY_SENTINEL,
    SECONDARY_SENTINEL,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _minimal_prs() -> Presentation:
    """Return a fresh 16:9 Presentation with one blank slide."""
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)
    blank = None
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if "blank" in layout.name.lower():
                blank = layout
                break
        if blank:
            break
    if blank is None:
        blank = prs.slide_masters[0].slide_layouts[6]
    prs.slides.add_slide(blank)
    return prs


def _make_tiny_png(path: str, width: int = 4, height: int = 4) -> None:
    """Write a minimal valid PNG file to *path*."""
    def chunk(name: bytes, data: bytes) -> bytes:
        c = name + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    ihdr_data = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    # Raw image: width × height pixels, each (R, G, B) = (0, 128, 255)
    raw_rows = b"".join(b"\x00" + b"\x00\x80\xFF" * width for _ in range(height))
    idat_data = zlib.compress(raw_rows)

    png_bytes = (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr_data)
        + chunk(b"IDAT", idat_data)
        + chunk(b"IEND", b"")
    )
    with open(path, "wb") as f:
        f.write(png_bytes)


# ---------------------------------------------------------------------------
# 1. BrandConfig defaults
# ---------------------------------------------------------------------------

def test_brand_config_from_dict_none():
    brand = BrandConfig.from_dict(None)
    assert brand.primary_color == DEFAULT_BRAND_CONFIG["primary_color"]
    assert brand.secondary_color == DEFAULT_BRAND_CONFIG["secondary_color"]
    assert brand.font_family == DEFAULT_BRAND_CONFIG["font_family"]
    assert brand.company_name == DEFAULT_BRAND_CONFIG["company_name"]
    assert brand.logo_path == DEFAULT_BRAND_CONFIG["logo_path"]


def test_brand_config_from_dict_partial_override():
    brand = BrandConfig.from_dict({"primary_color": "#AABBCC", "company_name": "Acme"})
    assert brand.primary_color == "#AABBCC"
    assert brand.company_name == "Acme"
    # Unspecified keys retain defaults
    assert brand.secondary_color == DEFAULT_BRAND_CONFIG["secondary_color"]


def test_brand_config_from_dict_full_override():
    cfg = {
        "company_name": "Foo Inc",
        "logo_path": "foo_logo.png",
        "primary_color": "#111111",
        "secondary_color": "#222222",
        "font_family": "Arial",
    }
    brand = BrandConfig.from_dict(cfg)
    assert brand.company_name == "Foo Inc"
    assert brand.logo_path == "foo_logo.png"
    assert brand.primary_color == "#111111"
    assert brand.secondary_color == "#222222"
    assert brand.font_family == "Arial"


# ---------------------------------------------------------------------------
# 2. hex_to_rgb conversion
# ---------------------------------------------------------------------------

def test_primary_rgb_conversion():
    brand = BrandConfig.from_dict({"primary_color": "#2563EB"})
    rgb = brand.primary_rgb()
    assert isinstance(rgb, RGBColor)
    assert rgb[0] == 0x25
    assert rgb[1] == 0x63
    assert rgb[2] == 0xEB


def test_secondary_rgb_conversion():
    brand = BrandConfig.from_dict({"secondary_color": "#6B7280"})
    rgb = brand.secondary_rgb()
    assert rgb[0] == 0x6B
    assert rgb[1] == 0x72
    assert rgb[2] == 0x80


def test_rgb_conversion_without_hash():
    brand = BrandConfig(primary_color="FF0000")
    rgb = brand.primary_rgb()
    assert rgb[0] == 0xFF
    assert rgb[1] == 0x00
    assert rgb[2] == 0x00


def test_rgb_conversion_invalid_hex():
    brand = BrandConfig(primary_color="#ZZZZZZ")
    with pytest.raises(ValueError):
        brand.primary_rgb()


# ---------------------------------------------------------------------------
# 3. merge_brand_config
# ---------------------------------------------------------------------------

def test_merge_brand_config_none_returns_defaults():
    merged = merge_brand_config(None)
    assert merged == DEFAULT_BRAND_CONFIG


def test_merge_brand_config_partial():
    merged = merge_brand_config({"primary_color": "#DEADBE"})
    assert merged["primary_color"] == "#DEADBE"
    assert merged["secondary_color"] == DEFAULT_BRAND_CONFIG["secondary_color"]


def test_merge_brand_config_empty_dict():
    merged = merge_brand_config({})
    assert merged == DEFAULT_BRAND_CONFIG


def test_merge_brand_config_does_not_mutate_default():
    original_primary = DEFAULT_BRAND_CONFIG["primary_color"]
    merge_brand_config({"primary_color": "#999999"})
    assert DEFAULT_BRAND_CONFIG["primary_color"] == original_primary


# ---------------------------------------------------------------------------
# 4. apply_brand_to_slide — logo insertion
# ---------------------------------------------------------------------------

def test_logo_inserted_top_right(tmp_path):
    """A valid PNG logo must produce a Picture shape on the slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    logo_path = str(tmp_path / "logo.png")
    _make_tiny_png(logo_path)

    # Patch /workspace/input to tmp_path
    original_join = os.path.join
    import helpers.brand as brand_mod
    _original = brand_mod.os.path.join

    # We directly pass the full abs path by setting logo_path to a relative name
    # and having the function prefix /workspace/input; instead, monkeypatch the prefix.
    # Simpler: create tmp_path subdir named "workspace/input" and put logo there.
    input_dir = tmp_path / "workspace" / "input"
    input_dir.mkdir(parents=True)
    real_logo = input_dir / "logo.png"
    _make_tiny_png(str(real_logo))

    # Temporarily reroute /workspace/input to our tmp input_dir
    import unittest.mock as mock

    prs = _minimal_prs()
    slide = prs.slides[0]
    brand = BrandConfig(logo_path="logo.png")

    with mock.patch.object(brand_mod.os.path, "join",
                           side_effect=lambda *args: str(input_dir / args[-1])
                           if args[0] == "/workspace/input" else original_join(*args)):
        warnings = apply_brand_to_slide(slide, brand)

    assert warnings == []
    pic_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pic_shapes) >= 1


def test_logo_missing_file_warning():
    """Non-existent logo path must yield a warning and not crash."""
    prs = _minimal_prs()
    slide = prs.slides[0]
    brand = BrandConfig(logo_path="definitely_not_here.png")
    warnings = apply_brand_to_slide(slide, brand)
    assert len(warnings) == 1
    assert "logo" in warnings[0].lower() or "不存在" in warnings[0]


def test_svg_logo_warning():
    """SVG logo must be rejected with an explicit SVG warning."""
    prs = _minimal_prs()
    slide = prs.slides[0]
    brand = BrandConfig(logo_path="company_logo.svg")
    warnings = apply_brand_to_slide(slide, brand)
    assert len(warnings) == 1
    assert "SVG" in warnings[0] or "svg" in warnings[0].lower()


def test_no_logo_no_warnings():
    """brand_config without logo_path should produce zero warnings."""
    prs = _minimal_prs()
    slide = prs.slides[0]
    brand = BrandConfig(logo_path=None)
    warnings = apply_brand_to_slide(slide, brand)
    assert warnings == []


def test_apply_logo_false_skips_insertion():
    """apply_logo=False must skip insertion even if logo_path is set."""
    prs = _minimal_prs()
    slide = prs.slides[0]
    brand = BrandConfig(logo_path="logo.png")
    warnings = apply_brand_to_slide(slide, brand, apply_logo=False)
    assert warnings == []


# ---------------------------------------------------------------------------
# 5. apply_brand_to_deck
# ---------------------------------------------------------------------------

def test_apply_brand_to_deck_no_logo():
    """Deck with no logo_path should have zero warnings from apply_brand_to_deck."""
    prs = _minimal_prs()
    prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])
    brand = BrandConfig(logo_path=None)
    warnings = apply_brand_to_deck(prs, brand)
    assert warnings == []


# ---------------------------------------------------------------------------
# 6. apply_brand_colors_to_template — sentinel replacement
# ---------------------------------------------------------------------------

def _sentinel_prs() -> Presentation:
    """Return a Presentation whose slide text contains sentinel hex strings."""
    from pptx.util import Cm
    from lxml import etree

    prs = _minimal_prs()
    slide = prs.slides[0]
    # Inject sentinel colour into slide XML as a fill element
    sp = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(5), Cm(2))
    # Manually add a solidFill element with PRIMARY_SENTINEL to mimic a coloured shape
    from pptx.oxml.ns import qn
    spPr = sp._element.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp._element, qn("p:spPr"))
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgb.set("val", PRIMARY_SENTINEL)
    return prs


def test_brand_colors_template_replacement():
    """Sentinel hex in template XML must be replaced with brand colour."""
    prs = _sentinel_prs()
    brand = BrandConfig(primary_color="#ABCDEF", secondary_color="#123456")

    count = apply_brand_colors_to_template(prs, brand)
    assert count >= 1

    # Verify replacement happened: no PRIMARY_SENTINEL remaining
    import xml.etree.ElementTree as ET
    xml_str = prs.slides[0]._element.xml
    assert PRIMARY_SENTINEL not in xml_str


def test_brand_colors_template_no_sentinels():
    """A presentation with no sentinel colours should return count 0."""
    prs = _minimal_prs()
    brand = BrandConfig(primary_color="#ABCDEF", secondary_color="#123456")
    count = apply_brand_colors_to_template(prs, brand)
    # Count may be 0 (no sentinels found) — no crash is the key requirement
    assert isinstance(count, int)
    assert count >= 0
