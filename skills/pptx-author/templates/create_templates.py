"""
create_templates.py — Bootstrap script to generate the three pptx template files.

Run once (or whenever template design changes):
    python create_templates.py

Creates:
    briefing.pptx   — dark-cover briefing / reporting template
    analysis.pptx   — light analytical / data-heavy template
    proposal.pptx   — split-layout proposal template

Sentinel colours used throughout:
    primary:    #003399  (replaced at runtime by apply_brand_colors_to_template)
    secondary:  #CC6600  (same)

Templates contain 10 slides each covering all 9 SLIDE_LAYOUTS + a spare,
so visual designers can use them as references when adjusting coordinates.
"""

from __future__ import annotations

import os
import sys

# Allow import from parent package when run as a script
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

PRIMARY = RGBColor(0x00, 0x33, 0x99)    # sentinel primary
SECONDARY = RGBColor(0xCC, 0x66, 0x00)  # sentinel secondary
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF3, 0xF4, 0xF6)
MID_GREY = RGBColor(0x6B, 0x72, 0x80)


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_layout(prs: Presentation):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if "blank" in layout.name.lower():
                return layout
    return prs.slide_masters[0].slide_layouts[6]


def _add_rect(slide, left, top, width, height, rgb: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    return shape


def _add_text(slide, text, left, top, width, height,
              font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def _make_briefing() -> Presentation:
    """Briefing template: dark-header slides, professional reporting style."""
    prs = _new_prs()
    bl = _blank_layout(prs)

    for i in range(10):
        slide = prs.slides.add_slide(bl)
        if i == 0:
            # Cover slide
            _add_rect(slide, 0, 0, SLIDE_W, int(SLIDE_H * 0.45), PRIMARY)
            _add_text(slide, "{{company}} 汇报标题", Cm(1.8), Cm(2.0),
                      SLIDE_W - Cm(3.6), Cm(3.5), font_size=32, bold=True, color=WHITE)
            _add_text(slide, "{{period}}  ·  {{department}}", Cm(1.8),
                      int(SLIDE_H * 0.45) + Cm(0.6), SLIDE_W - Cm(3.6), Cm(2.0),
                      font_size=16, color=MID_GREY)
        elif i == 1:
            # Section slide
            _add_rect(slide, 0, 0, Cm(0.6), SLIDE_H, PRIMARY)
            _add_text(slide, "01", Cm(2.0), SLIDE_H * 0.28, SLIDE_W - Cm(3.0),
                      Cm(3.0), font_size=48, bold=True, color=PRIMARY)
            _add_text(slide, "章节标题", Cm(2.0), SLIDE_H * 0.28 + Cm(3.0),
                      SLIDE_W - Cm(3.0), Cm(3.0), font_size=24, color=SECONDARY)
        elif i == 9:
            # End slide
            _add_rect(slide, 0, int(SLIDE_H * 0.92), SLIDE_W,
                      int(SLIDE_H * 0.08), PRIMARY)
            _add_text(slide, "谢谢", Cm(1.5), SLIDE_H * 0.30,
                      SLIDE_W - Cm(3.0), SLIDE_H * 0.25,
                      font_size=40, bold=True, color=PRIMARY,
                      align=PP_ALIGN.CENTER)
            _add_text(slide, "{{company}}", Cm(1.5),
                      int(SLIDE_H * 0.92) + Cm(0.15), SLIDE_W / 2, Cm(1.2),
                      font_size=11, color=WHITE)
        else:
            # Generic content slides
            _add_rect(slide, 0, 0, SLIDE_W, Cm(1.2), PRIMARY)
            _add_text(slide, f"幻灯片 {i + 1} 标题", Cm(1.5), Cm(1.4),
                      SLIDE_W - Cm(3.0), Cm(2.0), font_size=24, bold=True,
                      color=PRIMARY)
            _add_text(slide, "内容区域", Cm(1.5), Cm(4.0),
                      SLIDE_W - Cm(3.0), SLIDE_H - Cm(5.5),
                      font_size=14, color=MID_GREY)
    return prs


def _make_analysis() -> Presentation:
    """Analysis template: light background, large chart space."""
    prs = _new_prs()
    bl = _blank_layout(prs)

    for i in range(10):
        slide = prs.slides.add_slide(bl)
        # Light-grey top bar
        _add_rect(slide, 0, 0, SLIDE_W, Cm(1.0), LIGHT_GREY)
        if i == 0:
            _add_text(slide, "{{company}} 分析报告", Cm(1.5), Cm(2.5),
                      SLIDE_W - Cm(3.0), Cm(3.0), font_size=32, bold=True,
                      color=PRIMARY)
            _add_text(slide, "分析师: {{analyst}}  日期: {{date}}  数据集: {{dataset_name}}",
                      Cm(1.5), Cm(6.0), SLIDE_W - Cm(3.0), Cm(2.0),
                      font_size=14, color=MID_GREY)
        elif i == 9:
            _add_text(slide, "结论", Cm(1.5), Cm(2.5), SLIDE_W - Cm(3.0), Cm(2.5),
                      font_size=32, bold=True, color=PRIMARY)
            _add_rect(slide, 0, int(SLIDE_H * 0.92), SLIDE_W,
                      int(SLIDE_H * 0.08), SECONDARY)
        else:
            _add_text(slide, f"分析页 {i}", Cm(1.5), Cm(1.2),
                      SLIDE_W - Cm(3.0), Cm(2.0), font_size=22, bold=True,
                      color=PRIMARY)
            # Large chart placeholder area
            _add_rect(slide, Cm(1.5), Cm(3.8), SLIDE_W - Cm(3.0),
                      SLIDE_H - Cm(5.5), LIGHT_GREY)
            _add_text(slide, "[图表区域]", Cm(1.5), Cm(3.8) + (SLIDE_H - Cm(5.5)) / 2,
                      SLIDE_W - Cm(3.0), Cm(1.5), font_size=14, color=MID_GREY,
                      align=PP_ALIGN.CENTER)
    return prs


def _make_proposal() -> Presentation:
    """Proposal template: split left-stripe / right-content layout."""
    prs = _new_prs()
    bl = _blank_layout(prs)
    STRIPE_W = int(SLIDE_W * 0.33)

    for i in range(10):
        slide = prs.slides.add_slide(bl)
        _add_rect(slide, 0, 0, STRIPE_W, SLIDE_H, PRIMARY)
        if i == 0:
            _add_text(slide, "{{project_name}}", Cm(1.0), SLIDE_H * 0.25,
                      STRIPE_W - Cm(1.5), SLIDE_H * 0.5,
                      font_size=24, bold=True, color=WHITE)
            _add_text(slide, "客户: {{client}}", STRIPE_W + Cm(1.0), SLIDE_H * 0.30,
                      SLIDE_W - STRIPE_W - Cm(1.5), Cm(2.0),
                      font_size=18, color=MID_GREY)
            _add_text(slide, "版本: {{version}}", STRIPE_W + Cm(1.0),
                      SLIDE_H * 0.30 + Cm(2.5),
                      SLIDE_W - STRIPE_W - Cm(1.5), Cm(2.0),
                      font_size=14, color=MID_GREY)
        elif i == 9:
            _add_text(slide, "下一步", Cm(1.0), SLIDE_H * 0.35,
                      STRIPE_W - Cm(1.5), SLIDE_H * 0.3,
                      font_size=22, bold=True, color=WHITE)
            _add_text(slide, "立即开始", STRIPE_W + Cm(1.0), SLIDE_H * 0.35,
                      SLIDE_W - STRIPE_W - Cm(1.5), Cm(3.0),
                      font_size=28, bold=True, color=PRIMARY)
        else:
            _add_text(slide, f"提案页 {i}", Cm(1.0), SLIDE_H * 0.15,
                      STRIPE_W - Cm(1.5), Cm(3.0),
                      font_size=18, bold=True, color=WHITE)
            _add_text(slide, "内容区域", STRIPE_W + Cm(1.0), Cm(2.0),
                      SLIDE_W - STRIPE_W - Cm(1.5), SLIDE_H - Cm(4.0),
                      font_size=14, color=MID_GREY)
    return prs


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    templates = {
        "briefing.pptx": _make_briefing(),
        "analysis.pptx": _make_analysis(),
        "proposal.pptx": _make_proposal(),
    }
    for fname, prs in templates.items():
        path = os.path.join(out_dir, fname)
        prs.save(path)
        print(f"Created: {path}")


if __name__ == "__main__":
    main()
