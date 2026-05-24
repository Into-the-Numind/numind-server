"""
layout_renderers.py — One render function per SLIDE_LAYOUTS entry.

All renderers share the same signature::

    render_*(slide, slide_def, brand, style) -> None

Design decision: absolute coordinate positioning (Cm-based) rather than
placeholder index access.  Reasons:
  - pptx template placeholder indices are unstable across different templates.
  - Coordinate positioning lets render_* functions be reused across all
    three templates (briefing / analysis / proposal).
  - Brand colours are applied directly to new text runs (not patched after).

All renderers use slide_w / slide_h obtained from the prs object so that
they work correctly with any slide size (16:9 assumed by DeckBuilder defaults).
"""

from __future__ import annotations

import os
from typing import Optional, TYPE_CHECKING

from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

if TYPE_CHECKING:
    from pptx.slide import Slide
    from .brand import BrandConfig


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

_DEFAULT_STYLE: dict = {
    "title_font_size_pt": 28,
    "body_font_size_pt": 16,
    "bullet_font_size_pt": 14,
    "background_color": "#FFFFFF",
}


def _get_style(style: Optional[dict]) -> dict:
    s = dict(_DEFAULT_STYLE)
    if style:
        s.update(style)
    return s


def _slide_dims(slide: "Slide"):
    """Return (slide_width_emu, slide_height_emu) from the parent Presentation.

    Supports python-pptx >=0.6 (package.presentation) and >=1.0
    (package.presentation_part.presentation).
    """
    pkg = slide.part.package
    if hasattr(pkg, "presentation"):
        prs = pkg.presentation
    else:
        prs = pkg.presentation_part.presentation
    return prs.slide_width, prs.slide_height


def _primary_rgb(brand: Optional["BrandConfig"]) -> Optional[RGBColor]:
    return brand.primary_rgb() if brand else None


def _secondary_rgb(brand: Optional["BrandConfig"]) -> Optional[RGBColor]:
    return brand.secondary_rgb() if brand else None


def _font_name(brand: Optional["BrandConfig"], fallback: str = "Calibri") -> str:
    return brand.font_family if brand else fallback


def _add_title_box(slide: "Slide", title: str, brand, style: dict,
                   top_cm: float = 1.2, height_cm: float = 2.5,
                   margin_cm: float = 1.5) -> None:
    """Add a standard title textbox at the top of the slide."""
    slide_w, _ = _slide_dims(slide)
    box = slide.shapes.add_textbox(
        Cm(margin_cm), Cm(top_cm),
        slide_w - Cm(margin_cm * 2), Cm(height_cm),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(style["title_font_size_pt"])
    run.font.bold = True
    run.font.name = _font_name(brand)
    rgb = _primary_rgb(brand)
    if rgb:
        run.font.color.rgb = rgb


def _add_slide_background(slide: "Slide", hex_color: str) -> None:
    """Fill the slide background with a solid colour."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    h = hex_color.lstrip("#")
    fill.fore_color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _write_notes(slide: "Slide", notes_text: Optional[str]) -> None:
    """Write speaker notes to the notes slide if notes_text is non-empty."""
    if not notes_text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = notes_text
    except Exception:
        pass  # notes_slide creation can fail silently; non-fatal


# ---------------------------------------------------------------------------
# Slide layout renderers
# ---------------------------------------------------------------------------

def render_cover(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                 style: Optional[dict]) -> None:
    """
    Cover slide: full-width primary-colour header band, large title, subtitle,
    optional company name, and logo if configured.

    Layout:
        [  primary-colour rectangle (top 40% of slide)  ]
        [  title (large, white)                          ]
        [  subtitle / date / dept (below rect, dark)     ]
        [  company name (bottom right)                   ]
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)

    # ---- Background colour band (primary colour fills top portion) ----
    primary = _primary_rgb(brand)
    band_h = slide_h * 0.42

    rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        0, 0, slide_w, int(band_h),
    )
    rect.line.fill.background()  # no border
    fill = rect.fill
    fill.solid()
    if primary:
        fill.fore_color.rgb = primary
    else:
        fill.fore_color.rgb = RGBColor(0x1F, 0x29, 0x37)

    # ---- Main title (white, inside band) ----
    title_text = slide_def.get("title", "")
    title_box = slide.shapes.add_textbox(
        Cm(1.8), Cm(2.0), slide_w - Cm(3.6), int(band_h) - Cm(1.5),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(max(s["title_font_size_pt"], 32))
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = _font_name(brand)

    # ---- Subtitle (below band) ----
    subtitle_text = slide_def.get("subtitle", "")
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(
            Cm(1.8), int(band_h) + Cm(0.6), slide_w - Cm(3.6), Cm(3.0),
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle_text
        run.font.size = Pt(s["body_font_size_pt"])
        run.font.name = _font_name(brand)
        secondary = _secondary_rgb(brand)
        if secondary:
            run.font.color.rgb = secondary

    # ---- Company name (bottom left) ----
    company = (brand.company_name if brand else "") or slide_def.get("company_name", "")
    if company:
        co_box = slide.shapes.add_textbox(
            Cm(1.8), slide_h - Cm(1.8), slide_w / 2, Cm(1.4),
        )
        tf = co_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = company
        run.font.size = Pt(11)
        run.font.name = _font_name(brand)
        if primary:
            run.font.color.rgb = primary

    _write_notes(slide, slide_def.get("notes"))


def render_section(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                   style: Optional[dict]) -> None:
    """
    Section / transition slide: centred chapter number (large) + chapter title.

    Layout:
        [              (vertical space)                 ]
        [   Chapter N  (large, primary colour)          ]
        [   Section title (medium)                      ]
        [              (vertical space)                 ]
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)
    primary = _primary_rgb(brand)

    # Optional: left-side accent strip (thin vertical bar)
    accent_rect = slide.shapes.add_shape(
        1, 0, 0, Cm(0.6), slide_h,
    )
    accent_rect.line.fill.background()
    fill = accent_rect.fill
    fill.solid()
    if primary:
        fill.fore_color.rgb = primary
    else:
        fill.fore_color.rgb = RGBColor(0x1F, 0x29, 0x37)

    # Subtitle field doubles as chapter label (e.g. "Chapter 1" or "01")
    chapter_label = slide_def.get("subtitle", "")
    if chapter_label:
        label_box = slide.shapes.add_textbox(
            Cm(2.0), slide_h * 0.28, slide_w - Cm(3.0), Cm(3.0),
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = chapter_label
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.name = _font_name(brand)
        if primary:
            run.font.color.rgb = primary

    # Section title
    title_box = slide.shapes.add_textbox(
        Cm(2.0), slide_h * 0.28 + Cm(3.0), slide_w - Cm(3.0), Cm(4.0),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_def.get("title", "")
    run.font.size = Pt(s["title_font_size_pt"])
    run.font.bold = False
    run.font.name = _font_name(brand)
    secondary = _secondary_rgb(brand)
    if secondary:
        run.font.color.rgb = secondary

    _write_notes(slide, slide_def.get("notes"))


def render_title_body(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                      style: Optional[dict]) -> None:
    """
    Title + body text layout.  Body text is a single paragraph (supports \\n splitting).
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)

    _add_title_box(slide, slide_def.get("title", ""), brand, s)

    body_text = slide_def.get("body", "")
    if body_text:
        body_box = slide.shapes.add_textbox(
            Cm(1.5), Cm(4.0), slide_w - Cm(3.0), slide_h - Cm(5.5),
        )
        tf = body_box.text_frame
        tf.word_wrap = True

        paragraphs = body_text.split("\n")
        for i, para in enumerate(paragraphs):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            run = p.add_run()
            run.text = para
            run.font.size = Pt(s["body_font_size_pt"])
            run.font.name = _font_name(brand)

    _write_notes(slide, slide_def.get("notes"))


def render_title_bullets(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                          style: Optional[dict]) -> None:
    """
    Title + bullet point list (max 6 bullets).

    Each bullet is prefixed with a bullet character (•).
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)

    _add_title_box(slide, slide_def.get("title", ""), brand, s)

    bullet_points: list = slide_def.get("bullet_points", [])[:6]
    if bullet_points:
        bullet_box = slide.shapes.add_textbox(
            Cm(1.5), Cm(4.0), slide_w - Cm(3.0), slide_h - Cm(5.5),
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"•  {point}"
            run.font.size = Pt(s["bullet_font_size_pt"])
            run.font.name = _font_name(brand)

    _write_notes(slide, slide_def.get("notes"))


def render_title_table(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                        style: Optional[dict]) -> None:
    """
    Title + table layout.

    table dict schema::

        {
            "headers": ["Col A", "Col B", ...],
            "rows":    [["r1c1", "r1c2", ...], ...],
        }

    Header row uses primary_color as background with white text.
    Data rows alternate light-grey (#F3F4F6) and white.
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)
    primary = _primary_rgb(brand)

    _add_title_box(slide, slide_def.get("title", ""), brand, s)

    table_def = slide_def.get("table")
    if not table_def:
        _write_notes(slide, slide_def.get("notes"))
        return

    headers: list = table_def.get("headers", [])
    rows: list = table_def.get("rows", [])
    if not headers and not rows:
        _write_notes(slide, slide_def.get("notes"))
        return

    cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    total_rows = (1 if headers else 0) + len(rows)

    left = Cm(1.5)
    top = Cm(4.2)
    width = slide_w - Cm(3.0)
    height = min(Cm(1.0) * total_rows, slide_h - Cm(5.5))

    table = slide.shapes.add_table(total_rows, cols, left, top, width, height).table
    table.horz_banding = True

    font_sz = Pt(max(s["body_font_size_pt"] - 2, 10))
    font_name = _font_name(brand)

    row_offset = 0
    if headers:
        for ci, hdr in enumerate(headers[:cols]):
            cell = table.cell(0, ci)
            cell.text = hdr
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = hdr
            # Clear default paragraph text first
            tf.paragraphs[0].runs[0].text = hdr if tf.paragraphs[0].runs else ""
            run.font.size = font_sz
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = font_name
            # Header background
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            from pptx.oxml import parse_xml
            from pptx.oxml.ns import nsmap
            solidFill = parse_xml(
                f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{(primary or RGBColor(0x1F, 0x29, 0x37)).replace(" ", "").lstrip("(").rstrip(")")}"/>'
                f'</a:solidFill>'
            ) if False else None  # simplified — use direct attribute set
            # Set header background via fill
            if primary:
                from pptx.oxml.ns import qn as _qn
                spPr = tc.get_or_add_tcPr()
                _set_cell_bg(cell, primary)
        row_offset = 1

    for ri, row_data in enumerate(rows):
        for ci, cell_val in enumerate(row_data[:cols]):
            cell = table.cell(ri + row_offset, ci)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            # Clear existing then add run
            for run in p.runs:
                run.text = ""
            run = p.add_run()
            run.text = str(cell_val)
            run.font.size = font_sz
            run.font.name = font_name

    _write_notes(slide, slide_def.get("notes"))


def _set_cell_bg(cell, rgb: RGBColor) -> None:
    """Set table cell solid background colour via XML manipulation."""
    from lxml import etree
    from pptx.oxml.ns import qn as _qn

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Remove any existing fill
    for old in tcPr.findall(_qn("a:solidFill")):
        tcPr.remove(old)

    solid_fill = etree.SubElement(tcPr, _qn("a:solidFill"))
    srgb = etree.SubElement(solid_fill, _qn("a:srgbClr"))
    srgb.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def render_title_chart(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                        style: Optional[dict]) -> None:
    """
    Title + chart layout.  The chart is rendered via matplotlib and embedded
    as a PNG picture occupying the lower ~75% of the slide.
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)

    _add_title_box(slide, slide_def.get("title", ""), brand, s)

    chart_def = slide_def.get("chart")
    if not chart_def:
        _write_notes(slide, slide_def.get("notes"))
        return

    from .chart import render_chart_to_image, insert_chart_image

    size = chart_def.get("size", {})
    left = float(size.get("left", 0.05))
    top = float(size.get("top", 0.25))
    width = float(size.get("width", 0.90))
    height = float(size.get("height", 0.68))

    primary_hex = brand.primary_color if brand else None
    secondary_hex = brand.secondary_color if brand else None

    try:
        img_stream = render_chart_to_image(
            chart_def,
            primary_color=primary_hex,
            secondary_color=secondary_hex,
        )
        insert_chart_image(slide, img_stream, left, top, width, height)
    except Exception as exc:
        # Non-fatal: add a warning textbox instead of crashing
        warn_box = slide.shapes.add_textbox(Cm(1.5), Cm(4.0), slide_w - Cm(3.0), Cm(2.0))
        warn_box.text_frame.paragraphs[0].text = f"[图表渲染失败: {exc}]"

    _write_notes(slide, slide_def.get("notes"))


def render_title_image(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                        style: Optional[dict]) -> None:
    """
    Title + image layout.

    image dict schema::

        {
            "path": "relative/to/input/dir.jpg",
            "left":    0.1,    # fraction of slide width (default 0.1)
            "top":     0.25,   # fraction of slide height (default 0.25)
            "width":   0.80,   # fraction of slide width  (default 0.80)
            "caption": "图片说明文字",
        }
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)

    _add_title_box(slide, slide_def.get("title", ""), brand, s)

    img_def = slide_def.get("image")
    if not img_def or not img_def.get("path"):
        _write_notes(slide, slide_def.get("notes"))
        return

    img_path = os.path.join("/workspace/input", img_def["path"])
    left_frac = float(img_def.get("left", 0.1))
    top_frac = float(img_def.get("top", 0.25))
    width_frac = float(img_def.get("width", 0.80))

    left_emu = int(slide_w * left_frac)
    top_emu = int(slide_h * top_frac)
    width_emu = int(slide_w * width_frac)

    warnings: list[str] = []
    try:
        from PIL import Image as PILImage
        with PILImage.open(img_path) as pil_img:
            orig_w, orig_h = pil_img.size
        if orig_w > 0:
            height_emu = int(width_emu * orig_h / orig_w)
        else:
            height_emu = int(slide_h * 0.5)
        slide.shapes.add_picture(img_path, left_emu, top_emu, width_emu, height_emu)
    except FileNotFoundError:
        warn_box = slide.shapes.add_textbox(Cm(1.5), Cm(4.0), slide_w - Cm(3.0), Cm(1.5))
        warn_box.text_frame.paragraphs[0].text = f"[图片文件不存在: {img_def['path']}]"
    except Exception as exc:
        warn_box = slide.shapes.add_textbox(Cm(1.5), Cm(4.0), slide_w - Cm(3.0), Cm(1.5))
        warn_box.text_frame.paragraphs[0].text = f"[图片插入失败: {exc}]"

    # Caption
    caption = img_def.get("caption", "")
    if caption:
        cap_box = slide.shapes.add_textbox(
            left_emu, top_emu + int(slide_h * 0.6), width_emu, Cm(1.2),
        )
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = caption
        run.font.size = Pt(11)
        run.font.name = _font_name(brand)
        secondary = _secondary_rgb(brand)
        if secondary:
            run.font.color.rgb = secondary

    _write_notes(slide, slide_def.get("notes"))


def render_two_column(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
                       style: Optional[dict]) -> None:
    """
    Two-column layout: left side body/bullets, right side chart or image.

    Uses ``body`` / ``bullet_points`` for the left column and
    ``chart`` or ``image`` for the right column.
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)

    _add_title_box(slide, slide_def.get("title", ""), brand, s)

    col_top = Cm(4.0)
    col_h = slide_h - Cm(5.5)
    col_w = (slide_w - Cm(3.5)) / 2
    left_left = Cm(1.5)
    right_left = left_left + col_w + Cm(0.5)

    # ---- Left column: body or bullets ----
    left_content = slide_def.get("body") or ""
    bullets = slide_def.get("bullet_points", [])

    if bullets:
        left_box = slide.shapes.add_textbox(left_left, col_top, col_w, col_h)
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(bullets[:6]):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = f"•  {pt}"
            run.font.size = Pt(s["bullet_font_size_pt"])
            run.font.name = _font_name(brand)
    elif left_content:
        left_box = slide.shapes.add_textbox(left_left, col_top, col_w, col_h)
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, para in enumerate(left_content.split("\n")):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            run = p.add_run()
            run.text = para
            run.font.size = Pt(s["body_font_size_pt"])
            run.font.name = _font_name(brand)

    # ---- Right column: chart or image ----
    chart_def = slide_def.get("chart")
    img_def = slide_def.get("image")

    right_left_frac = right_left / slide_w
    top_frac = col_top / slide_h
    width_frac = col_w / slide_w
    height_frac = col_h / slide_h

    if chart_def:
        from .chart import render_chart_to_image, insert_chart_image
        try:
            img_stream = render_chart_to_image(
                chart_def,
                primary_color=brand.primary_color if brand else None,
            )
            insert_chart_image(slide, img_stream, right_left_frac, top_frac,
                               width_frac, height_frac)
        except Exception as exc:
            wb = slide.shapes.add_textbox(int(right_left), int(col_top), int(col_w), Cm(1.5))
            wb.text_frame.paragraphs[0].text = f"[图表渲染失败: {exc}]"

    elif img_def and img_def.get("path"):
        img_path = os.path.join("/workspace/input", img_def["path"])
        try:
            from PIL import Image as PILImage
            with PILImage.open(img_path) as pil_img:
                ow, oh = pil_img.size
            w_emu = int(col_w)
            h_emu = int(col_w * oh / ow) if ow > 0 else int(col_h)
            slide.shapes.add_picture(img_path, int(right_left), int(col_top), w_emu, h_emu)
        except Exception as exc:
            wb = slide.shapes.add_textbox(int(right_left), int(col_top), int(col_w), Cm(1.5))
            wb.text_frame.paragraphs[0].text = f"[右列内容: {exc}]"

    _write_notes(slide, slide_def.get("notes"))


def render_end(slide: "Slide", slide_def: dict, brand: Optional["BrandConfig"],
               style: Optional[dict]) -> None:
    """
    End / thank-you slide: centred headline, optional sub-text, company name,
    and logo.  Uses the same primary-colour band as cover for visual bookending.
    """
    s = _get_style(style)
    slide_w, slide_h = _slide_dims(slide)
    primary = _primary_rgb(brand)
    secondary = _secondary_rgb(brand)

    # Full-width accent strip at bottom 8%
    accent_h = int(slide_h * 0.08)
    rect = slide.shapes.add_shape(1, 0, slide_h - accent_h, slide_w, accent_h)
    rect.line.fill.background()
    fill = rect.fill
    fill.solid()
    if primary:
        fill.fore_color.rgb = primary
    else:
        fill.fore_color.rgb = RGBColor(0x1F, 0x29, 0x37)

    # Main thank-you text (centred)
    main_text = slide_def.get("title", "谢谢")
    title_box = slide.shapes.add_textbox(
        Cm(1.5), slide_h * 0.30, slide_w - Cm(3.0), slide_h * 0.25,
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = main_text
    run.font.size = Pt(max(s["title_font_size_pt"], 36))
    run.font.bold = True
    run.font.name = _font_name(brand)
    if primary:
        run.font.color.rgb = primary

    # Subtitle (e.g. "Q&A" or contact info)
    subtitle = slide_def.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Cm(1.5), slide_h * 0.55, slide_w - Cm(3.0), slide_h * 0.18,
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(s["body_font_size_pt"])
        run.font.name = _font_name(brand)
        if secondary:
            run.font.color.rgb = secondary

    # Company name (bottom-left of accent strip)
    company = (brand.company_name if brand else "") or slide_def.get("company_name", "")
    if company:
        co_box = slide.shapes.add_textbox(
            Cm(1.5), slide_h - accent_h + Cm(0.15), slide_w / 2, int(accent_h * 0.7),
        )
        tf = co_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = company
        run.font.size = Pt(11)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = _font_name(brand)

    _write_notes(slide, slide_def.get("notes"))
