"""
chart.py — matplotlib-based chart rendering helpers for pptx-author.

Design decision: charts are rendered to PNG (in-memory BytesIO) via matplotlib
and embedded as pictures in the slide.  This provides:
  - Universal compatibility (PowerPoint / WPS / Google Slides / LibreOffice)
  - Brand-colour integration (primary_color applied to bars/lines)
  - No dependency on python-pptx's native chart XML (which has poor compat)

IMPORTANT: matplotlib.use("Agg") must be called BEFORE importing pyplot.
This module sets the backend at import time so callers don't need to care.
"""

from __future__ import annotations

import matplotlib  # noqa: E402  — backend must be set before pyplot import
matplotlib.use("Agg")  # headless / no-display mode for sandbox

import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
from io import BytesIO
from typing import Optional

from pptx.util import Cm


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _has_cjk_font() -> bool:
    """
    Return True if a CJK-capable font is available in the current matplotlib
    font cache.  Checks for WenQuanYi Zen Hei (sandbox pre-installed), SimHei,
    Noto Sans CJK, and Microsoft YaHei.
    """
    known_cjk = {"wenquanyi zen hei", "simhei", "noto sans cjk sc", "microsoft yahei"}
    for f in fm.fontManager.ttflist:
        if f.name.lower() in known_cjk:
            return True
        if any(kw in f.fname.lower() for kw in ("wqy", "simhei", "notosanscjk")):
            return True
    return False


def _get_cjk_font_prop() -> Optional[fm.FontProperties]:
    """
    Return a FontProperties for CJK rendering if available, else None.
    Used for axis labels and titles that may contain Chinese characters.
    """
    cjk_names = [
        "WenQuanYi Zen Hei",
        "Noto Sans CJK SC",
        "SimHei",
        "Microsoft YaHei",
    ]
    for name in cjk_names:
        try:
            prop = fm.FontProperties(family=name)
            # Test that matplotlib can resolve this font
            fm.findfont(prop, fallback_to_default=False)
            return prop
        except Exception:
            continue
    return None


def _build_colors(primary_color: Optional[str], n: int) -> Optional[list]:
    """
    Build a list of *n* matplotlib colours derived from *primary_color*.
    The first colour is the brand primary; subsequent colours are generated
    by shifting lightness through the Blues colormap.

    Returns None if primary_color is absent (matplotlib uses its own cycle).
    """
    if not primary_color:
        return None
    try:
        h = primary_color.lstrip("#")
        r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255
        primary_rgb = (r, g, b)
    except (ValueError, AttributeError):
        return None

    if n == 1:
        return [primary_rgb]

    # Blend from primary toward Blues colormap for additional series
    extras = plt.cm.Blues(np.linspace(0.35, 0.75, n - 1))
    return [primary_rgb] + [tuple(c) for c in extras]


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def render_chart_to_image(
    chart_def: dict,
    width_px: int = 900,
    height_px: int = 500,
    dpi: int = 150,
    primary_color: Optional[str] = None,
    secondary_color: Optional[str] = None,
) -> BytesIO:
    """
    Render a chart definition to a PNG image stored in a BytesIO buffer.

    Supports chart types: ``"bar"``, ``"line"``, ``"pie"``, ``"scatter"``.

    Parameters
    ----------
    chart_def:
        Dictionary conforming to ``InvokeParams.slides[n].chart``.
        Required keys: ``chart_type``, ``data`` (with ``categories`` and
        ``series`` lists), optional ``title``.
    width_px, height_px, dpi:
        Image dimensions.  Default gives a 900×500 px image at 150 dpi.
    primary_color:
        CSS hex string (e.g. ``"#2563EB"``) used as the first bar/line colour.
        Subsequent series derive from a tinted version of this colour.
    secondary_color:
        Currently unused; reserved for future two-colour chart styles.

    Returns
    -------
    BytesIO
        Seeked-to-position-0 PNG stream ready to pass to ``add_picture``.

    Raises
    ------
    ValueError
        If ``chart_type`` is not one of the supported types.
    KeyError
        If required keys are missing from ``chart_def`` or ``data``.
    """
    chart_type = chart_def.get("chart_type", "bar").lower()
    if chart_type not in ("bar", "line", "pie", "scatter"):
        raise ValueError(
            f"Unsupported chart_type: {chart_type!r}. "
            "Supported: bar, line, pie, scatter."
        )

    data = chart_def["data"]
    categories: list = data["categories"]
    series: list = data["series"]

    figsize = (width_px / dpi, height_px / dpi)
    fig, ax = plt.subplots(figsize=figsize, dpi=dpi)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#F9FAFB")  # subtle off-white plot area

    colors = _build_colors(primary_color, len(series))
    cjk_prop = _get_cjk_font_prop()

    # ------------------------------------------------------------------ bar
    if chart_type == "bar":
        x = np.arange(len(categories))
        bar_w = 0.8 / max(len(series), 1)
        for i, s in enumerate(series):
            color = colors[i] if colors else None
            offset = (i - (len(series) - 1) / 2) * bar_w
            ax.bar(x + offset, s["values"], bar_w, label=s["name"], color=color)
        ax.set_xticks(x)
        if cjk_prop:
            ax.set_xticklabels(categories, rotation=20 if len(categories) > 5 else 0,
                               fontproperties=cjk_prop, fontsize=9)
        else:
            ax.set_xticklabels(categories, rotation=20 if len(categories) > 5 else 0,
                               fontsize=9)
        ax.yaxis.grid(True, linestyle="--", alpha=0.6)
        ax.set_axisbelow(True)

    # ----------------------------------------------------------------- line
    elif chart_type == "line":
        for i, s in enumerate(series):
            color = colors[i] if colors else None
            ax.plot(categories, s["values"], marker="o", linewidth=2,
                    label=s["name"], color=color)
        ax.yaxis.grid(True, linestyle="--", alpha=0.6)
        ax.set_axisbelow(True)
        if cjk_prop and len(categories) > 0:
            ax.set_xticklabels(ax.get_xticklabels(), fontproperties=cjk_prop, fontsize=9)

    # ------------------------------------------------------------------ pie
    elif chart_type == "pie":
        values = series[0]["values"] if series else []
        pie_colors = (plt.cm.Set2(np.linspace(0, 1, len(values)))
                      if not colors else colors[:len(values)])
        wedge_props = {"linewidth": 0.5, "edgecolor": "white"}
        if cjk_prop:
            ax.pie(values, labels=categories, autopct="%1.1f%%",
                   colors=pie_colors, wedgeprops=wedge_props,
                   textprops={"fontproperties": cjk_prop})
        else:
            ax.pie(values, labels=categories, autopct="%1.1f%%",
                   colors=pie_colors, wedgeprops=wedge_props)
        ax.axis("equal")

    # --------------------------------------------------------------- scatter
    elif chart_type == "scatter":
        for i, s in enumerate(series):
            color = colors[i] if colors else None
            try:
                x_vals = [float(c) for c in categories]
            except (ValueError, TypeError):
                # Fallback: use integer indices when categories aren't numeric
                x_vals = list(range(len(s["values"])))
            ax.scatter(x_vals, s["values"], label=s["name"], color=color, alpha=0.7)
        ax.yaxis.grid(True, linestyle="--", alpha=0.6)
        ax.set_axisbelow(True)

    # -------------------------------------------------------------- Title / legend
    chart_title = chart_def.get("title", "")
    if chart_title:
        if cjk_prop:
            ax.set_title(chart_title, fontsize=13, fontproperties=cjk_prop)
        else:
            ax.set_title(chart_title, fontsize=13)

    if len(series) > 1:
        if cjk_prop:
            ax.legend(prop=cjk_prop, fontsize=9)
        else:
            ax.legend(fontsize=9)

    ax.tick_params(labelsize=9)
    plt.tight_layout(pad=0.8)

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def insert_chart_image(
    slide,
    image_stream: BytesIO,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """
    Embed a PNG image stream (from render_chart_to_image) into *slide*
    at a position expressed as fractions of the slide dimensions.

    Parameters
    ----------
    slide:
        A ``pptx.slide.Slide`` object.
    image_stream:
        BytesIO containing a PNG (seeked to position 0).
    left, top, width, height:
        Fractions of slide width/height in [0.0, 1.0].
    """
    pkg = slide.part.package
    prs = pkg.presentation if hasattr(pkg, "presentation") else pkg.presentation_part.presentation
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    left_emu = int(slide_w * left)
    top_emu = int(slide_h * top)
    width_emu = int(slide_w * width)
    height_emu = int(slide_h * height)

    slide.shapes.add_picture(image_stream, left_emu, top_emu, width_emu, height_emu)
