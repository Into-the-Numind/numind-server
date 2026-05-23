"""
deck.py — DeckBuilder: slide lifecycle manager for pptx-author.

DeckBuilder is the single entry point for callers:

1. Instantiate with optional template_path, brand config, and style config.
2. Call add_slide(slide_def) for each slide in order.
3. Call apply_brand() to insert logos across all slides.
4. Call save(output_path) to write the .pptx file.

Design decisions:
- 16:9 dimensions set explicitly on Presentation() — python-pptx defaults to 4:3.
- Layout dispatch uses a mapping dict; unknown layout names fall back to "title-body".
- Blank layout is located by name ("Blank") rather than by index — more robust
  across different template files.
- apply_brand() is separate from add_slide() so logos are inserted after all
  content is placed (avoids z-order issues with full-bleed shapes in cover slides).
"""

from __future__ import annotations

import os
from typing import Optional

from pptx import Presentation
from pptx.util import Cm, Emu

from .brand import BrandConfig, apply_brand_to_deck, apply_brand_colors_to_template
from .layout_renderers import (
    render_cover,
    render_section,
    render_title_body,
    render_title_bullets,
    render_title_table,
    render_title_chart,
    render_title_image,
    render_two_column,
    render_end,
)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# 16:9 widescreen dimensions (33.87 cm × 19.05 cm)
SLIDE_WIDTH_CM = 33.87
SLIDE_HEIGHT_CM = 19.05

# Mapping from layout name → renderer function
_LAYOUT_DISPATCH: dict = {
    "cover": render_cover,
    "section": render_section,
    "title-body": render_title_body,
    "title-bullets": render_title_bullets,
    "title-table": render_title_table,
    "title-chart": render_title_chart,
    "title-image": render_title_image,
    "two-column": render_two_column,
    "blank": None,   # blank — no renderer; add_slide returns the raw slide
    "end": render_end,
}

_TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "..", "templates")


class DeckBuilder:
    """
    Builder for a PowerPoint Presentation (.pptx).

    Parameters
    ----------
    template_path:
        Optional path to an existing .pptx file used as the base.
        When provided, the existing slides from the template are kept and
        brand sentinel colours are replaced.  New slides are appended.
        Use the ``template`` short-name ("briefing" | "analysis" | "proposal")
        or a full absolute path.
    brand:
        BrandConfig instance.  If None, no brand colours/fonts are applied.
    style_config:
        Dict of style overrides (title_font_size_pt, body_font_size_pt, etc.).
        See layout_renderers._DEFAULT_STYLE for available keys.
    """

    def __init__(
        self,
        template_path: Optional[str] = None,
        brand: Optional[BrandConfig] = None,
        style_config: Optional[dict] = None,
    ) -> None:
        self._brand = brand
        self._style = style_config or {}
        self._warnings: list[str] = []

        if template_path:
            resolved = self._resolve_template(template_path)
            self._prs = Presentation(resolved)
            # Sentinel colour replacement for pre-built templates
            if brand:
                apply_brand_colors_to_template(self._prs, brand)
        else:
            self._prs = Presentation()

        # Always enforce 16:9 dimensions (override template's aspect ratio too)
        slide_w_cm = self._style.pop("slide_width_cm", SLIDE_WIDTH_CM)
        slide_h_cm = self._style.pop("slide_height_cm", SLIDE_HEIGHT_CM)
        self._prs.slide_width = Cm(slide_w_cm)
        self._prs.slide_height = Cm(slide_h_cm)

        self._blank_layout = self._find_blank_layout()

    # -----------------------------------------------------------------------
    # Internal helpers
    # -----------------------------------------------------------------------

    def _resolve_template(self, name_or_path: str) -> str:
        """Resolve a short template name or absolute path to a filesystem path."""
        if os.path.isabs(name_or_path) and os.path.exists(name_or_path):
            return name_or_path
        # Short name → look in templates/
        candidates = [
            os.path.join(_TEMPLATE_DIR, f"{name_or_path}.pptx"),
            os.path.join(_TEMPLATE_DIR, name_or_path),
        ]
        for cand in candidates:
            if os.path.exists(cand):
                return cand
        raise FileNotFoundError(
            f"Template not found: {name_or_path!r}. "
            f"Tried: {candidates}"
        )

    def _find_blank_layout(self):
        """
        Locate the Blank slide layout in the Presentation's slide master.

        Searches by name (case-insensitive) to be robust across templates.
        Falls back to index 6 (standard position in most MS Office templates)
        and then to index 0 if nothing matches.
        """
        for master in self._prs.slide_masters:
            for layout in master.slide_layouts:
                if "blank" in layout.name.lower():
                    return layout
        # Fallback: try index 6 (common blank position)
        try:
            return self._prs.slide_masters[0].slide_layouts[6]
        except IndexError:
            return self._prs.slide_masters[0].slide_layouts[0]

    # -----------------------------------------------------------------------
    # Public API
    # -----------------------------------------------------------------------

    def add_slide(self, slide_def: dict) -> "DeckBuilder":
        """
        Add one slide to the Presentation.

        Parameters
        ----------
        slide_def:
            A dict conforming to the InvokeParams.slides[] schema.
            Required: ``layout`` (str from SLIDE_LAYOUTS).
            Optional: ``title``, ``subtitle``, ``body``, ``bullet_points``,
                      ``table``, ``chart``, ``image``, ``notes``.

        Returns
        -------
        DeckBuilder
            self (for optional method chaining).
        """
        layout_name = slide_def.get("layout", "title-body").lower()
        renderer = _LAYOUT_DISPATCH.get(layout_name)

        # Add blank slide as base
        slide = self._prs.slides.add_slide(self._blank_layout)

        if renderer is not None:
            try:
                renderer(slide, slide_def, self._brand, self._style)
            except Exception as exc:
                self._warnings.append(
                    f"slide {len(self._prs.slides)}: layout={layout_name!r} "
                    f"render error — {exc}"
                )
        # "blank" layout: slide is returned empty (caller's intent)

        return self

    def apply_brand(self) -> "DeckBuilder":
        """
        Insert the brand logo into every slide.

        This is called after all slides have been added so that logo z-order
        is above all content shapes.

        Returns
        -------
        DeckBuilder
            self (for optional method chaining).
        """
        if self._brand:
            warnings = apply_brand_to_deck(self._prs, self._brand)
            self._warnings.extend(warnings)
        return self

    def save(self, output_path: str) -> None:
        """
        Write the Presentation to *output_path*.

        Parameters
        ----------
        output_path:
            Absolute filesystem path (e.g. ``"/output/report.pptx"``).
        """
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self._prs.save(output_path)

    @property
    def warnings(self) -> list[str]:
        """Non-fatal warning messages accumulated during build."""
        return list(self._warnings)

    @property
    def slide_count(self) -> int:
        """Number of slides currently in the Presentation."""
        return len(self._prs.slides)
