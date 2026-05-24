"""
brand.py — Brand identity configuration and application helpers.

Provides:
  - DEFAULT_BRAND_CONFIG: neutral defaults (T18 decision: generic blank)
  - BrandConfig dataclass with hex-to-RGBColor conversion helpers
  - merge_brand_config: deep-merge user-supplied config onto defaults
  - apply_brand_to_slide: insert logo + (optionally) apply color/font to a slide
  - apply_brand_to_deck: batch apply brand across all slides in a Presentation
  - apply_brand_colors_to_template: sentinel-hex replacement for template-based decks
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from io import BytesIO
from typing import Optional

from pptx.util import Cm
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------------------
# Decision T18: Default brand config — generic/neutral (non-Numind branding).
# Users override any key via the brand_config parameter in invoke_skill.
# ---------------------------------------------------------------------------
DEFAULT_BRAND_CONFIG: dict = {
    "company_name": "",                # empty — caller fills in
    "logo_path": None,                 # None → no logo rendered
    "primary_color": "#1F2937",        # neutral dark-grey (Tailwind gray-800)
    "secondary_color": "#6B7280",      # neutral mid-grey  (Tailwind gray-500)
    "font_family": "Noto Sans CJK SC", # CJK-capable; sandbox alias → wqy-zenhei
}

# Sentinel hex values embedded in the pptx template files.
# These values act as placeholders for brand-colour replacement.
# They have no visual meaning on their own — chosen to be rarely-used colours.
PRIMARY_SENTINEL = "003399"    # deep blue placeholder for primary
SECONDARY_SENTINEL = "CC6600"  # amber placeholder for secondary


def merge_brand_config(user_config: Optional[dict]) -> dict:
    """
    Deep-merge *user_config* onto DEFAULT_BRAND_CONFIG.

    Only keys present in user_config are overwritten; the rest retain their
    defaults.  Returns a new dict (does not mutate DEFAULT_BRAND_CONFIG).

    Example::

        cfg = merge_brand_config({"primary_color": "#2563EB", "company_name": "Acme"})
        # cfg["secondary_color"] is still "#6B7280" (the default)
    """
    merged = dict(DEFAULT_BRAND_CONFIG)
    if user_config:
        for key, value in user_config.items():
            if value is not None or key in ("logo_path", "company_name"):
                merged[key] = value
    return merged


@dataclass
class BrandConfig:
    """
    Typed wrapper around the merged brand configuration dict.

    Attributes mirror the keys in DEFAULT_BRAND_CONFIG / InvokeParams.brand_config.
    """

    company_name: str = ""
    logo_path: Optional[str] = None
    primary_color: str = DEFAULT_BRAND_CONFIG["primary_color"]
    secondary_color: str = DEFAULT_BRAND_CONFIG["secondary_color"]
    font_family: str = DEFAULT_BRAND_CONFIG["font_family"]

    # ------------------------------------------------------------------
    # Colour helpers
    # ------------------------------------------------------------------

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert a CSS hex colour string (with or without leading #) to RGBColor."""
        h = hex_color.lstrip("#")
        if len(h) != 6:
            raise ValueError(f"Invalid hex colour: {hex_color!r}")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return RGBColor(r, g, b)

    def primary_rgb(self) -> RGBColor:
        """Return the primary colour as a pptx RGBColor."""
        return self._hex_to_rgb(self.primary_color)

    def secondary_rgb(self) -> RGBColor:
        """Return the secondary colour as a pptx RGBColor."""
        return self._hex_to_rgb(self.secondary_color)

    # ------------------------------------------------------------------
    # Factory
    # ------------------------------------------------------------------

    @classmethod
    def from_dict(cls, cfg: dict) -> "BrandConfig":
        """Build a BrandConfig from a (possibly partial) config dict."""
        merged = merge_brand_config(cfg)
        return cls(
            company_name=merged.get("company_name", ""),
            logo_path=merged.get("logo_path"),
            primary_color=merged.get("primary_color", DEFAULT_BRAND_CONFIG["primary_color"]),
            secondary_color=merged.get("secondary_color", DEFAULT_BRAND_CONFIG["secondary_color"]),
            font_family=merged.get("font_family", DEFAULT_BRAND_CONFIG["font_family"]),
        )


# ---------------------------------------------------------------------------
# Logo insertion
# ---------------------------------------------------------------------------

def apply_brand_to_slide(
    slide,
    brand: BrandConfig,
    apply_logo: bool = True,
    logo_position: str = "top-right",
    logo_width_cm: float = 2.5,
) -> list[str]:
    """
    Apply brand identity to a single slide.

    Currently handles logo insertion (PNG/JPG) while preserving aspect ratio.
    Colour / font application is handled at render time inside layout_renderers.py
    (colours are set directly on new text runs, not patched after the fact).

    Parameters
    ----------
    slide:
        A ``pptx.slide.Slide`` object.
    brand:
        A BrandConfig instance.
    apply_logo:
        Whether to attempt logo insertion.  Set False for blank/section slides
        where the logo would overlap content.
    logo_position:
        One of ``"top-right"`` (default), ``"bottom-right"``, ``"bottom-left"``.
    logo_width_cm:
        Rendered width of the logo in centimetres.  Height is computed to
        preserve the original aspect ratio.

    Returns
    -------
    list[str]
        Non-fatal warning messages (e.g. missing file, unsupported format).
    """
    warnings: list[str] = []

    if not apply_logo or not brand.logo_path:
        return warnings

    # SVG is not supported by Pillow — check early and warn.
    if brand.logo_path.lower().endswith(".svg"):
        warnings.append(
            "logo 格式不支持 SVG，请提供 PNG 或 JPG 格式的 logo 文件，已跳过 logo 插入。"
        )
        return warnings

    logo_abs = os.path.join("/workspace/input", brand.logo_path)

    try:
        from PIL import Image as PILImage

        with PILImage.open(logo_abs) as img:
            orig_w, orig_h = img.size

        if orig_w == 0:
            warnings.append(f"logo 文件宽度为 0，跳过 logo 插入: {brand.logo_path}")
            return warnings

        logo_h_cm = logo_width_cm * (orig_h / orig_w)

        pkg = slide.part.package
        prs = pkg.presentation if hasattr(pkg, "presentation") else pkg.presentation_part.presentation
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        margin = Cm(0.5)
        logo_w_emu = Cm(logo_width_cm)
        logo_h_emu = Cm(logo_h_cm)

        if logo_position == "top-right":
            left = slide_w - logo_w_emu - margin
            top = margin
        elif logo_position == "bottom-right":
            left = slide_w - logo_w_emu - margin
            top = slide_h - logo_h_emu - margin
        elif logo_position == "bottom-left":
            left = margin
            top = slide_h - logo_h_emu - margin
        else:
            left = slide_w - logo_w_emu - margin
            top = margin

        slide.shapes.add_picture(logo_abs, left, top, logo_w_emu, logo_h_emu)

    except FileNotFoundError:
        warnings.append(f"logo 文件不存在: {brand.logo_path}，已跳过 logo 插入。")
    except Exception as exc:
        warnings.append(f"logo 插入失败 ({exc})，已跳过。")

    return warnings


def apply_brand_to_deck(
    prs,
    brand: BrandConfig,
    logo_position: str = "top-right",
    logo_width_cm: float = 2.5,
) -> list[str]:
    """
    Batch-apply brand identity (logo only) to every slide in *prs*.

    Returns the aggregated list of warning strings from all slides.
    """
    all_warnings: list[str] = []
    for slide in prs.slides:
        warnings = apply_brand_to_slide(
            slide,
            brand,
            apply_logo=True,
            logo_position=logo_position,
            logo_width_cm=logo_width_cm,
        )
        all_warnings.extend(warnings)
    return all_warnings


# ---------------------------------------------------------------------------
# Template colour replacement (sentinel-hex substitution)
# ---------------------------------------------------------------------------

def _replace_color_in_element(element, old_hex: str, new_hex: str) -> int:
    """
    Recursively walk *element* XML tree replacing all ``val`` attribute values
    that equal *old_hex* with *new_hex*.

    Both hex strings must be 6-character uppercase/lowercase strings without
    a leading ``#``.  Returns the number of replacements made.
    """
    count = 0
    for elem in element.iter():
        if elem.get("val") == old_hex:
            elem.set("val", new_hex)
            count += 1
    return count


def apply_brand_colors_to_template(prs, brand: BrandConfig) -> int:
    """
    Replace sentinel placeholder colours in a template-loaded Presentation
    with the brand's primary / secondary colours.

    This is the mechanism used when DeckBuilder is initialised with a
    ``template_path`` (briefing / analysis / proposal .pptx).  The templates
    embed ``PRIMARY_SENTINEL`` and ``SECONDARY_SENTINEL`` as marker colours;
    this function swaps them for the caller's brand colours.

    Returns the total number of colour substitutions performed.
    """
    new_primary = brand.primary_color.lstrip("#")
    new_secondary = brand.secondary_color.lstrip("#")

    total = 0
    for slide in prs.slides:
        total += _replace_color_in_element(slide._element, PRIMARY_SENTINEL, new_primary)
        total += _replace_color_in_element(slide._element, SECONDARY_SENTINEL, new_secondary)

    for slide_master in prs.slide_masters:
        total += _replace_color_in_element(slide_master._element, PRIMARY_SENTINEL, new_primary)
        total += _replace_color_in_element(slide_master._element, SECONDARY_SENTINEL, new_secondary)

    return total
